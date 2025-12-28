import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import json
import os
import tempfile
from io import BytesIO

# --- 설정 및 상수 ---
CONFIG_FILE = 'config.json'
SAVED_TEMPLATE_PATH = 'saved_template.pptx'

st.set_page_config(
    page_title="AI 예배 PPT 생성기",
    page_icon="🕊️",
    layout="centered"
)

# 깔끔한 UI CSS
st.markdown("""
    <style>
    .main { padding-top: 2rem; }
    .stButton>button {
        width: 100%;
        border-radius: 8px;
        height: 3em;
        font-weight: bold;
    }
    .success-box {
        padding: 1rem;
        background-color: #f0fdf4;
        border: 1px solid #bbf7d0;
        border-radius: 0.5rem;
        color: #166534;
        margin-bottom: 1rem;
    }
    </style>
""", unsafe_allow_html=True)

# --- 0. 저장/불러오기 헬퍼 함수 ---

def load_api_key():
    """저장된 API 키 불러오기"""
    if os.path.exists(CONFIG_FILE):
        with open(CONFIG_FILE, "r") as f:
            return json.load(f).get("api_key", "")
    return ""

def save_api_key(key):
    """API 키 저장하기"""
    with open(CONFIG_FILE, "w") as f:
        json.dump({"api_key": key}, f)

def save_template(uploaded_file):
    """템플릿 파일 저장하기"""
    with open(SAVED_TEMPLATE_PATH, "wb") as f:
        f.write(uploaded_file.getbuffer())

# --- 1. 핵심 로직 (AI 및 PPT 처리) ---

def analyze_jubo_deep(image_file, key):
    genai.configure(api_key=key)
    model = genai.GenerativeModel('gemini-3-flash-preview') 
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp:
        tmp.write(image_file.getvalue())
        tmp_path = tmp.name

    try:
        sample_file = genai.upload_file(path=tmp_path)
        prompt = """
        이 주보 이미지에서 다음 정보를 찾아 JSON으로 출력해.
        값이 없으면 빈 문자열("")로 둬.
        
        1. sermon_title: 설교 제목
        2. preacher: 설교자 이름 (직분 포함)
        3. prayer_person: 대표 기도자 이름
        4. bible_ref: 성경 본문 위치 (예: 요한복음 3:16)
        5. bible_text: 위 bible_ref에 해당하는 실제 성경 말씀 내용을 '개역개정' 버전으로 찾아서 전체 작성해줘.
        6. hymn_list: 찬송가 제목들을 순서대로 리스트에 담아줘.
        """
        response = model.generate_content([sample_file, prompt])
        text = response.text.replace("```json", "").replace("```", "").strip()
        return json.loads(text)
    except Exception as e:
        st.error(f"분석 중 오류가 발생했습니다: {e}")
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def fill_ppt_text(template_path_or_file, data):
    # 파일 경로(문자열)인지 업로드된 파일 객체인지 확인하여 로드
    prs = Presentation(template_path_or_file)
    
    replacements = {
        "{{설교제목}}": data.get('sermon_title', ''),
        "{{설교자}}": data.get('preacher', ''),
        "{{기도자}}": data.get('prayer_person', ''),
        "{{성경본문}}": data.get('bible_ref', ''),
        "{{말씀내용}}": data.get('bible_text', '')
    }
    
    hymns = data.get('hymn_list', [])
    for i, hymn in enumerate(hymns):
        replacements[f"{{{{찬송{i+1}}}}}"] = hymn
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for k, v in replacements.items():
                        if k in run.text:
                            safe_value = str(v) if v is not None else ""
                            run.text = run.text.replace(k, safe_value)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- 2. 화면 UI 구성 ---

st.title("🕊️ AI 예배 PPT 생성기")
st.markdown("주보 사진만 올리면, 저장된 템플릿에 내용을 자동으로 채워줍니다.")

# 사이드바: 설정 (API 키 저장 기능 추가)
with st.sidebar:
    st.header("⚙️ 설정")
    
    # 저장된 키 불러오기
    saved_key = load_api_key()
    api_key_input = st.text_input("Google API Key", value=saved_key, type="password", placeholder="API 키를 입력하세요")
    
    # 키가 변경되면 저장
    if api_key_input != saved_key:
        save_api_key(api_key_input)
        st.success("API 키가 이 컴퓨터에 저장되었습니다!")
        # 즉시 반영을 위해 변수 업데이트
        saved_key = api_key_input

    st.divider()
    with st.expander("❓ 템플릿 가이드"):
        st.markdown("""
        템플릿 텍스트 상자 규칙:
        - `{{설교제목}}`, `{{설교자}}`
        - `{{기도자}}`, `{{성경본문}}`
        - `{{말씀내용}}`
        - `{{찬송1}}`, `{{찬송2}}`
        """)

# 메인 기능 영역
if not api_key_input:
    st.warning("👈 왼쪽 사이드바에 Google API 키를 먼저 입력해주세요.")
else:
    # STEP 1: 파일 준비
    st.subheader("1. 파일 준비")
    
    col1, col2 = st.columns(2)
    
    # 1-1. 주보 이미지 (항상 새로 업로드)
    with col1:
        st.markdown("**📸 주보 이미지**")
        jubo_img = st.file_uploader("주보 업로드", type=['png', 'jpg', 'jpeg'], label_visibility="collapsed")

    # 1-2. 템플릿 (저장된 것 확인 또는 변경)
    with col2:
        st.markdown("**📂 PPT 템플릿**")
        
        template_exists = os.path.exists(SAVED_TEMPLATE_PATH)
        target_template = None
        
        if template_exists:
            st.success("✅ 저장된 템플릿 사용 중")
            # 템플릿 변경 옵션
            new_template = st.file_uploader("템플릿 변경하기 (선택)", type=['pptx'], label_visibility="collapsed")
            if new_template:
                save_template(new_template)
                st.toast("새로운 템플릿이 저장되었습니다!")
                target_template = new_template
            else:
                target_template = SAVED_TEMPLATE_PATH
        else:
            st.info("등록된 템플릿이 없습니다.")
            new_template = st.file_uploader("템플릿 최초 등록", type=['pptx'])
            if new_template:
                save_template(new_template)
                st.rerun() # 저장 후 새로고침하여 '사용 중' 상태로 전환

    # STEP 2: AI 분석 실행
    if jubo_img and target_template:
        st.divider()
        if st.button("주보 분석 시작 ✨", type="primary"):
            with st.spinner("주보를 읽고 성경 말씀을 찾는 중입니다..."):
                result = analyze_jubo_deep(jubo_img, api_key_input)
                if result:
                    st.session_state['ppt_data'] = result
                    st.rerun()
    elif not jubo_img:
        st.info("👆 주보 이미지를 올려주세요.")
    elif not target_template:
        st.warning("👆 PPT 템플릿을 등록해주세요.")

    # STEP 3: 결과 확인 및 다운로드
    if 'ppt_data' in st.session_state:
        st.divider()
        st.subheader("2. 내용 확인 및 수정")
        
        st.markdown('<div class="success-box">✅ AI 분석 완료! 내용을 확인하세요.</div>', unsafe_allow_html=True)
        
        d = st.session_state['ppt_data']
        
        with st.form("check_form"):
            c1, c2 = st.columns(2)
            with c1:
                new_title = st.text_input("설교 제목", value=d.get('sermon_title', ''))
                new_preacher = st.text_input("설교자", value=d.get('preacher', ''))
            with c2:
                new_prayer = st.text_input("기도자", value=d.get('prayer_person', ''))
                new_ref = st.text_input("성경 본문", value=d.get('bible_ref', ''))
            
            new_text = st.text_area("성경 말씀 내용 (AI 자동 생성)", value=d.get('bible_text', ''), height=150)
            hymn_str = st.text_input("찬송가 순서 (쉼표로 구분)", value=", ".join(d.get('hymn_list', [])))
            
            submitted = st.form_submit_button("이 내용으로 PPT 만들기 🎁", type="primary")
            
            if submitted:
                # 템플릿 파일이 경로(str)인지 파일객체인지 다시 확인 (새로고침 등으로 변수가 날아갈 수 있음)
                # 가장 확실한 방법: 저장된 파일 경로 사용 (위에서 저장했으므로)
                final_template_source = SAVED_TEMPLATE_PATH if os.path.exists(SAVED_TEMPLATE_PATH) else None
                
                if final_template_source:
                    final_data = {
                        "sermon_title": new_title,
                        "preacher": new_preacher,
                        "prayer_person": new_prayer,
                        "bible_ref": new_ref,
                        "bible_text": new_text,
                        "hymn_list": [h.strip() for h in hymn_str.split(',')]
                    }
                    
                    final_ppt = fill_ppt_text(final_template_source, final_data)
                    
                    st.session_state['final_ppt'] = final_ppt
                    st.session_state['final_file_name'] = f"{new_title}_예배.pptx"
                else:
                    st.error("템플릿 파일을 찾을 수 없습니다. 다시 업로드해주세요.")

        if 'final_ppt' in st.session_state:
            st.divider()
            st.balloons()
            st.success("작업 완료! 아래 버튼으로 다운로드하세요.")
            st.download_button(
                label="📥 완성된 PPT 다운로드",
                data=st.session_state['final_ppt'],
                file_name=st.session_state['final_file_name'],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
